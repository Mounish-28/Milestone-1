import sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    NAVY = RGBColor(15, 23, 42)
    BLUE = RGBColor(37, 99, 235)
    LIGHT_BG = RGBColor(248, 250, 252)
    CARD_BG = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 41, 59)
    MUTED_TEXT = RGBColor(100, 116, 139)
    WHITE = RGBColor(255, 255, 255)
    BORDER_COLOR = RGBColor(226, 232, 240)

    def add_header(slide, title_text, category_text="SHOPSENSE PLATFORM"):
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(1.1))
        tf = header_box.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(11)
        p0.font.bold = True
        p0.font.color.rgb = BLUE
        p0.font.name = "Calibri"
        
        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(26)
        p1.font.bold = True
        p1.font.color.rgb = NAVY
        p1.font.name = "Calibri"

    def set_bg(slide, color):
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
        bg.fill.solid()
        bg.fill.fore_color.rgb = color
        bg.line.fill.background()

    # ---------------------------------------------------
    # Slide 1: Title Slide (Dark Background)
    # ---------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    set_bg(slide1, NAVY)

    box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    p0.text = "🛍️ ShopSense Analytics Platform"
    p0.font.size = Pt(38)
    p0.font.bold = True
    p0.font.color.rgb = WHITE
    p0.font.name = "Calibri"

    p1 = tf.add_paragraph()
    p1.text = "Full-Stack Architecture, REST API Design & Live Frontend Integration"
    p1.font.size = Pt(20)
    p1.font.color.rgb = RGBColor(148, 163, 184)
    p1.space_before = Pt(14)
    p1.font.name = "Calibri"

    p2 = tf.add_paragraph()
    p2.text = "Milestone 1 Capstone Project  |  FastAPI (Python) + React 19 (Vite) + SQLite"
    p2.font.size = Pt(14)
    p2.font.bold = True
    p2.font.color.rgb = BLUE
    p2.space_before = Pt(28)
    p2.font.name = "Calibri"

    # ---------------------------------------------------
    # Slide 2: Executive Summary & Goals
    # ---------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    set_bg(slide2, LIGHT_BG)
    add_header(slide2, "Executive Summary & Objectives")

    cards = [
        ("🎯 Project Objective", "Build an enterprise multi-vendor e-commerce analytics platform offering interactive dashboards, catalog management, and partner metrics."),
        ("⚡ Core Architecture", "Modern decoupled architecture pairing a FastAPI Python backend with a React 19 + Vite frontend and SQLite ORM database."),
        ("🔌 Seamless Integration", "Full API connectivity featuring dynamic environment URLs, Vite dev proxy, CORS middleware, and live /health connection telemetry.")
    ]

    lefts = [Inches(0.8), Inches(4.8), Inches(8.8)]
    for i, (title, desc) in enumerate(cards):
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, lefts[i], Inches(1.8), Inches(3.7), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tb = slide2.shapes.add_textbox(lefts[i] + Inches(0.2), Inches(2.0), Inches(3.3), Inches(4.4))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(14)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(12)

    # ---------------------------------------------------
    # Slide 3: System Architecture
    # ---------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    set_bg(slide3, LIGHT_BG)
    add_header(slide3, "Full-Stack System Architecture")

    layers = [
        ("1. Frontend Tier", "React 19 + Vite", "Single-Page App built with React Router v7, Recharts data visualization, ThemeContext (Light/Dark mode), and Framer Motion micro-animations."),
        ("2. API & Service Layer", "Axios & Vite Proxy", "Centralized services/api.js module, dynamic VITE_API_URL environment configuration, and /api dev proxy routing to avoid CORS."),
        ("3. Backend Tier", "FastAPI (Python)", "High-performance async FastAPI backend with Pydantic request/response schemas, CORS middleware, and interactive OpenAPI /docs."),
        ("4. Database Tier", "SQLite & SQLAlchemy", "Relational database schema managing Vendors, Products, and Customers with SQLAlchemy ORM session lifecycle management.")
    ]

    tops = [Inches(1.7), Inches(3.0), Inches(4.3), Inches(5.6)]
    for i, (tag, subtitle, detail) in enumerate(layers):
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), tops[i], Inches(11.733), Inches(1.1))
        box.fill.solid()
        box.fill.fore_color.rgb = CARD_BG
        box.line.color.rgb = BORDER_COLOR

        tb = slide3.shapes.add_textbox(Inches(1.0), tops[i] + Inches(0.1), Inches(11.333), Inches(0.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f"{tag} — {subtitle}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p2 = tf.add_paragraph()
        p2.text = detail
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(4)

    # ---------------------------------------------------
    # Slide 4: Key Platform Features
    # ---------------------------------------------------
    slide4 = prs.slides.add_slide(blank_layout)
    set_bg(slide4, LIGHT_BG)
    add_header(slide4, "Core Platform Features")

    features = [
        ("📊 Analytics Dashboard", "Real-time metrics, interactive revenue trend charts, and latest transaction history."),
        ("📦 Product Inventory", "Searchable, filterable product table with live create/delete API integration."),
        ("🏬 Vendor Management", "Marketplace vendor onboarding, email directory, and product assignment counts."),
        ("👥 Customer Intelligence", "Profiles, geographic city filtering, phone contact lookup, and order statistics."),
        ("📑 Multi-Format Exports", "Instant client-side generation of PDF, Excel (.xlsx), and CSV analytics reports."),
        ("🌙 Dynamic Theme System", "Smooth Dark & Light mode toggle powered by central CSS design tokens.")
    ]

    col_lefts = [Inches(0.8), Inches(6.8)]
    row_tops = [Inches(1.8), Inches(3.6), Inches(5.4)]

    for idx, (ft_title, ft_desc) in enumerate(features):
        col = idx % 2
        row = idx // 2

        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_lefts[col], row_tops[row], Inches(5.7), Inches(1.5))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tb = slide4.shapes.add_textbox(col_lefts[col] + Inches(0.2), row_tops[row] + Inches(0.15), Inches(5.3), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = ft_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = NAVY

        p2 = tf.add_paragraph()
        p2.text = ft_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(4)

    # ---------------------------------------------------
    # Slide 5: Frontend-Backend Integration
    # ---------------------------------------------------
    slide5 = prs.slides.add_slide(blank_layout)
    set_bg(slide5, LIGHT_BG)
    add_header(slide5, "Frontend-Backend Connection Architecture")

    conn_steps = [
        ("1. Dynamic API Base URL", "Configured import.meta.env.VITE_API_URL || 'http://localhost:8000' in services/api.js."),
        ("2. Vite Dev Server Proxy", "Added proxy rule in vite.config.js mapping /api to http://127.0.0.1:8000."),
        ("3. FastAPI CORS Middleware", "Configured origins allowing http://localhost:5173, 3000, and 127.0.0.1."),
        ("4. Live Connection Telemetry", "Header polls /health every 15s displaying live '🟢 API Online' or '🔴 Offline Mode'.")
    ]

    for idx, (step_title, step_desc) in enumerate(conn_steps):
        top = Inches(1.8 + idx * 1.3)
        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.1))
        card.fill.solid()
        card.fill.fore_color.rgb = CARD_BG
        card.line.color.rgb = BORDER_COLOR

        tb = slide5.shapes.add_textbox(Inches(1.0), top + Inches(0.15), Inches(11.333), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = BLUE

        p2 = tf.add_paragraph()
        p2.text = step_desc
        p2.font.size = Pt(13)
        p2.font.color.rgb = DARK_TEXT
        p2.space_before = Pt(4)

    # ---------------------------------------------------
    # Slide 6: Production Deployment & Conclusion (Dark Background)
    # ---------------------------------------------------
    slide6 = prs.slides.add_slide(blank_layout)
    set_bg(slide6, NAVY)

    header_box = slide6.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.733), Inches(1.1))
    tf = header_box.text_frame
    tf.word_wrap = True
    
    p0 = tf.paragraphs[0]
    p0.text = "PRODUCTION DEPLOYMENT & SUMMARY".upper()
    p0.font.size = Pt(11)
    p0.font.bold = True
    p0.font.color.rgb = BLUE
    
    p1 = tf.add_paragraph()
    p1.text = "Deployment Readiness & Future Roadmap"
    p1.font.size = Pt(26)
    p1.font.bold = True
    p1.font.color.rgb = WHITE

    box = slide6.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.8))
    tf = box.text_frame
    tf.word_wrap = True

    items = [
        ("📦 Prepared Deployment Configs", "Created Dockerfile, Procfile, requirements.txt, and vercel.json for 1-click cloud publishing."),
        ("⚡ Production Build Verified", "Frontend build bundled into dist/ via Vite with zero errors (built in 16.84s)."),
        ("🔮 Future Roadmap", "Planned enhancements include JWT authentication, payment gateways, and predictive inventory analytics.")
    ]

    for title, desc in items:
        p = tf.add_paragraph()
        p.text = f"• {title}: "
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.space_before = Pt(18)

        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(203, 213, 225)

    output_path = r"d:\Milestone 1\ShopSense_Milestone1_Presentation.pptx"
    prs.save(output_path)
    print(f"Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_presentation()
